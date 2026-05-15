import csv, re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR

CSV = Path('/Users/nova/.openclaw/workspace/supportdev_april_gid472044289.csv')
OUT = Path('/Users/nova/.openclaw/workspace/Catch_up_SupportDev_April_2026_TH.pptx')

rows=list(csv.DictReader(open(CSV), fieldnames=None))
# find header manually because csv has blank top rows
raw=list(csv.reader(open(CSV)))
for i,r in enumerate(raw):
    if r and r[0]=='YEAR':
        header=r; data=raw[i+1:]; break
records=[]
for r in data:
    if len(r)>=20 and r[0]=='2026' and r[1].isdigit():
        records.append(dict(zip(header,r)))
april=[r for r in records if int(r['WEEK NO.']) in [14,15,16,17,18]]

def num(r,k):
    v=r[k].strip().replace(',','')
    return int(v) if v and v!='No Data' else 0

def pct(r,k):
    return float(r[k].replace('%','') or 0)

def dur_min(s):
    if not s or s=='No Data': return 0
    h=re.search(r'(\d+)\s*h',s); m=re.search(r'(\d+)\s*m',s)
    return (int(h.group(1)) if h else 0)*60 + (int(m.group(1)) if m else 0)

def fmt_hm(minutes):
    h=round(minutes/60,1)
    return f'{h:.1f} ชม.'

total=sum(num(r,'Total Ticket') for r in april)
amaze=sum(num(r,'AMAZE') for r in april)
phoenix=sum(num(r,'Phoenix') for r in april)
kd=sum(num(r,'KD') for r in april)
miss_fr=sum(num(r,'MISSED FIRST RESPONSE') for r in april)
miss_sla=sum(num(r,'MISSED SLA') for r in april)
resolved=sum(num(r,'Resolve Ticket') for r in april)
gchat=sum(num(r,'ggChat') for r in april)
avg_fr=sum(pct(r,'% First Response') for r in april)/len(april)
avg_sla=sum(pct(r,'% SLA') for r in april)/len(april)
wavg_resp=sum(dur_min(r['AVG. Response Time'])*num(r,'Total Ticket') for r in april)/total
wavg_res=sum(dur_min(r['AVG. Resolution Time'])*num(r,'Total Ticket') for r in april)/total
peak=max(april,key=lambda r:num(r,'Total Ticket'))
best_sla=max(april,key=lambda r:pct(r,'% SLA'))

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
NAVY=RGBColor(18,35,56); BLUE=RGBColor(44,123,229); CYAN=RGBColor(30,190,190); GREEN=RGBColor(39,174,96); ORANGE=RGBColor(242,153,74); RED=RGBColor(235,87,87); GRAY=RGBColor(95,111,123); LIGHT=RGBColor(244,247,251); WHITE=RGBColor(255,255,255)
FONT='Arial'

def add_title(slide,title,subtitle=None):
    tb=slide.shapes.add_textbox(Inches(.55), Inches(.35), Inches(12.2), Inches(.55)); p=tb.text_frame.paragraphs[0]; p.text=title; p.font.bold=True; p.font.size=Pt(28); p.font.color.rgb=NAVY; p.font.name=FONT
    if subtitle:
        tb=slide.shapes.add_textbox(Inches(.58), Inches(.9), Inches(12), Inches(.35)); p=tb.text_frame.paragraphs[0]; p.text=subtitle; p.font.size=Pt(12); p.font.color.rgb=GRAY; p.font.name=FONT

def bullet_box(slide, x,y,w,h, title, bullets, color=BLUE):
    box=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)); box.fill.solid(); box.fill.fore_color.rgb=WHITE; box.line.color.rgb=RGBColor(225,230,238)
    tb=slide.shapes.add_textbox(Inches(x+.25), Inches(y+.18), Inches(w-.5), Inches(.35)); p=tb.text_frame.paragraphs[0]; p.text=title; p.font.bold=True; p.font.size=Pt(15); p.font.color.rgb=color; p.font.name=FONT
    tf=slide.shapes.add_textbox(Inches(x+.25), Inches(y+.62), Inches(w-.45), Inches(h-.75)).text_frame
    tf.word_wrap=True
    for idx,b in enumerate(bullets):
        p=tf.paragraphs[0] if idx==0 else tf.add_paragraph(); p.text='• '+b; p.font.size=Pt(13); p.font.color.rgb=NAVY; p.font.name=FONT; p.space_after=Pt(7)

def metric(slide,x,y,w,h,label,value,sub='',color=BLUE):
    box=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)); box.fill.solid(); box.fill.fore_color.rgb=color; box.line.color.rgb=color
    tf=box.text_frame; tf.clear();
    p=tf.paragraphs[0]; p.text=value; p.alignment=PP_ALIGN.CENTER; p.font.bold=True; p.font.size=Pt(25); p.font.color.rgb=WHITE; p.font.name=FONT
    p=tf.add_paragraph(); p.text=label; p.alignment=PP_ALIGN.CENTER; p.font.size=Pt(12); p.font.color.rgb=WHITE; p.font.name=FONT
    if sub:
        p=tf.add_paragraph(); p.text=sub; p.alignment=PP_ALIGN.CENTER; p.font.size=Pt(9); p.font.color.rgb=WHITE; p.font.name=FONT

def bar_chart(slide,x,y,w,h, labels, values, color=BLUE):
    maxv=max(values) or 1
    for i,(lab,val) in enumerate(zip(labels,values)):
        yy=y+i*(h/len(values));
        slide.shapes.add_textbox(Inches(x), Inches(yy+.05), Inches(1.2), Inches(.25)).text_frame.paragraphs[0].text=lab
        bg=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x+1.25), Inches(yy+.05), Inches(w-2), Inches(.22)); bg.fill.solid(); bg.fill.fore_color.rgb=RGBColor(230,235,242); bg.line.fill.background()
        fg=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x+1.25), Inches(yy+.05), Inches((w-2)*val/maxv), Inches(.22)); fg.fill.solid(); fg.fill.fore_color.rgb=color; fg.line.fill.background()
        tb=slide.shapes.add_textbox(Inches(x+w-.65), Inches(yy), Inches(.6), Inches(.28)); p=tb.text_frame.paragraphs[0]; p.text=str(val); p.font.size=Pt(11); p.font.color.rgb=NAVY

# slide 1
s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=LIGHT
add_title(s,'Catch up SupportDev April 2026','สรุปภาพรวมจาก Support Performance Report | Weeks 14–18 (30 Mar–3 May)')
metric(s,.75,1.65,2.2,1.15,'Total Tickets',str(total),'รวม 5 สัปดาห์',BLUE)
metric(s,3.2,1.65,2.2,1.15,'Avg First Response',f'{avg_fr:.1f}%',f'Missed {miss_fr}',ORANGE)
metric(s,5.65,1.65,2.2,1.15,'Avg SLA',f'{avg_sla:.1f}%',f'Missed {miss_sla}',GREEN)
metric(s,8.1,1.65,2.2,1.15,'Resolved',str(resolved),f'Backlog ล่าสุด {num(april[-1],"Backlog")}',CYAN)
metric(s,10.55,1.65,2.0,1.15,'ggChat',str(gchat),'เคสจาก chat',NAVY)
bullet_box(s,.75,3.25,5.8,2.6,'Executive Summary',[f'ปริมาณงานเดือนเมษายนอยู่ที่ {total} tickets โดย peak ที่ Week {peak["WEEK NO."]} ({num(peak,"Total Ticket")} tickets)',f'SLA ภาพรวมค่อนข้างแข็งแรง เฉลี่ย {avg_sla:.1f}% และ Week 18 แตะ 100%',f'จุดที่ควรจับตาคือ First Response เฉลี่ย {avg_fr:.1f}% โดย Week 15 ต่ำสุด 74%'],BLUE)
bullet_box(s,6.85,3.25,5.7,2.6,'Lead Takeaway',['ปัญหาหลักไม่ใช่ resolution แต่เป็น response discipline ในบางสัปดาห์','Phoenix spike ใน Week 16 และ Week 18 ทำให้ workload แกว่ง','แนะนำทำ weekly guardrail สำหรับ missed first response + chat intake'],ORANGE)
# slide2 trend
s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=LIGHT; add_title(s,'Ticket Volume Trend','แนวโน้มปริมาณงานรายสัปดาห์ในช่วง April')
labels=[f'W{r["WEEK NO."]}' for r in april]; vals=[num(r,'Total Ticket') for r in april]
bar_chart(s,1.0,1.5,10.8,3.0,labels,vals,BLUE)
bullet_box(s,1.0,4.8,5.4,1.55,'Insight',[f'Week 16 สูงสุด {num(peak,"Total Ticket")} tickets จาก Phoenix {num(peak,"Phoenix")} tickets', 'หลัง Week 16 volume ลดลงชัดเจน แต่กลับขึ้นอีกใน Week 18'],BLUE)
bullet_box(s,6.8,4.8,5.0,1.55,'Risk',['การแกว่งของ volume สะท้อน dependency กับ event/spike เฉพาะระบบ','ควรแยก incident-like spike ออกจาก BAU เพื่ออ่าน capacity ให้แม่นขึ้น'],RED)
# slide3 mix
s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=LIGHT; add_title(s,'Workload Mix by System','สัดส่วน workload แยกตามระบบ')
metric(s,1,1.5,3,1.2,'AMAZE',str(amaze),f'{amaze/total*100:.1f}% ของทั้งหมด',BLUE)
metric(s,5.15,1.5,3,1.2,'Phoenix',str(phoenix),f'{phoenix/total*100:.1f}% ของทั้งหมด',ORANGE)
metric(s,9.3,1.5,3,1.2,'KD',str(kd),f'{kd/total*100:.1f}% ของทั้งหมด',GRAY)
bar_chart(s,1.2,3.25,10.2,1.9,['AMAZE','Phoenix','KD'],[amaze,phoenix,kd],CYAN)
bullet_box(s,1.2,5.35,10.2,1.0,'Interpretation',[f'AMAZE ยังเป็น workload หลัก ({amaze} tickets) แต่ Phoenix มี spike สำคัญใน Week 16/18 รวม {phoenix} tickets'],NAVY)
# slide4 SLA
s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=LIGHT; add_title(s,'SLA & Resolution Health','คุณภาพการปิดเคสและความเสี่ยง SLA')
metric(s,.9,1.45,2.7,1.15,'Avg SLA',f'{avg_sla:.1f}%','weekly average',GREEN)
metric(s,3.95,1.45,2.7,1.15,'Missed SLA',str(miss_sla),'รวมเมษายน',RED)
metric(s,7.0,1.45,2.7,1.15,'Avg Resolution',fmt_hm(wavg_res),'weighted by ticket',BLUE)
metric(s,10.05,1.45,2.4,1.15,'Resolved',str(resolved),f'{resolved/total*100:.1f}% of tickets',CYAN)
bar_chart(s,1,3.15,10.8,2.2,labels,[int(pct(r,'% SLA')) for r in april],GREEN)
bullet_box(s,1,5.75,10.8,.8,'Lead View',['Resolution health ดีขึ้นช่วงปลายเดือน: Week 17–18 missed SLA รวมแค่ 1 เคส และ Week 18 = 100%'],GREEN)
# slide5 response
s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=LIGHT; add_title(s,'First Response Focus','จุดที่ควรปรับเพื่อไม่ให้ SLA ดีแต่ customer wait สูง')
metric(s,1,1.45,3,1.15,'Avg First Response',f'{avg_fr:.1f}%','weekly average',ORANGE)
metric(s,4.45,1.45,3,1.15,'Missed First Response',str(miss_fr),'รวมเมษายน',RED)
metric(s,7.9,1.45,3,1.15,'Avg Response Time',fmt_hm(wavg_resp),'weighted by ticket',BLUE)
bar_chart(s,1.1,3.15,10.7,2.1,labels,[int(pct(r,'% First Response')) for r in april],ORANGE)
bullet_box(s,1.1,5.55,10.7,1.0,'Action Needed',['Week 15 ต่ำสุดที่ 74% และ missed 21 เคส — ควรทำ root cause แยก: queue coverage, handoff, holiday/weekend pattern, ggChat routing'],RED)
# slide6 recommendations
s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=LIGHT; add_title(s,'Recommended Next Actions','ข้อเสนอเพื่อ stabilize SupportDev cadence')
bullet_box(s,.8,1.35,3.85,4.5,'1) Response Guardrail',['ตั้ง daily threshold: missed first response > 5 ให้แจ้ง lead','ทำ queue sweep 2 รอบ/วันในช่วง volume สูง','แยก owner สำหรับ ggChat intake'],ORANGE)
bullet_box(s,4.85,1.35,3.85,4.5,'2) Spike Readiness',['Tag เคส Phoenix spike แยกจาก BAU','ทำ mini-RCA เฉพาะ Week 16 ว่ามาจาก issue/theme ไหน','เตรียม playbook เมื่อ Phoenix > 40 tickets/week'],BLUE)
bullet_box(s,8.9,1.35,3.65,4.5,'3) Reporting Hygiene',['เพิ่ม MoM summary ในชีต','แยก unresolved/backlog movement ให้ชัด','เพิ่ม top issue/theme เพื่อให้ management เห็น action ไม่ใช่แค่ตัวเลข'],GREEN)
# slide7 appendix
s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=LIGHT; add_title(s,'Appendix: Weekly Data','ข้อมูลรายสัปดาห์ที่ใช้ทำสรุป')
cols=['Week','Period','Tickets','AMAZE','Phoenix','FR%','SLA%','Resolved','ggChat']
table=s.shapes.add_table(len(april)+1,len(cols),Inches(.55),Inches(1.35),Inches(12.2),Inches(4.6)).table
for j,c in enumerate(cols):
    cell=table.cell(0,j); cell.text=c; cell.fill.solid(); cell.fill.fore_color.rgb=NAVY; cell.text_frame.paragraphs[0].font.color.rgb=WHITE; cell.text_frame.paragraphs[0].font.bold=True; cell.text_frame.paragraphs[0].font.size=Pt(10)
for i,r in enumerate(april,1):
    vals=[r['WEEK NO.'],f"{r['From Date']}–{r['To Date']}",r['Total Ticket'],r['AMAZE'],r['Phoenix'],r['% First Response'],r['% SLA'],r['Resolve Ticket'],r['ggChat']]
    for j,v in enumerate(vals):
        cell=table.cell(i,j); cell.text=str(v); cell.text_frame.paragraphs[0].font.size=Pt(9); cell.text_frame.paragraphs[0].font.color.rgb=NAVY
prs.save(OUT)
print(OUT)